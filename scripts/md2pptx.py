#!/usr/bin/env python3
"""Render a repo markdown doc to a PowerPoint deck with python-pptx only.

Companion to scripts/md2pdf.py for SME review: one slide per heading (h2-h4),
content greedy-packed with continuation slides, pipe tables as real pptx
tables (header repeated on splits), fenced code in monospace boxes (the §3.4
struct tree fits on one slide), inline **bold** / *italic* / `code` / [links]
as styled runs, and the slide's source text mirrored into speaker notes.
No transliteration needed — PowerPoint fonts cover the tree glyphs.

Decks are regenerable artifacts and git-ignored — never commit one:

    .venv/bin/python scripts/md2pptx.py                     # docs/ICD-NN.md -> docs/ICD-NN.pptx
    .venv/bin/python scripts/md2pptx.py other.md out.pptx
"""

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

SLIDE_W, SLIDE_H = 13.333, 7.5
CONTENT_TOP, CONTENT_BOTTOM, MARGIN = 1.22, 6.95, 0.55
BUDGET = CONTENT_BOTTOM - CONTENT_TOP

INK = RGBColor(0x20, 0x24, 0x2A)
ACCENT = RGBColor(0x1D, 0x3A, 0x5F)
DIM = RGBColor(0x77, 0x7F, 0x8A)
CODE_BG = RGBColor(0xF3, 0xF3, 0xF3)
HEADER_BG = RGBColor(0xE9, 0xE9, 0xE9)
BORDER = RGBColor(0x99, 0x99, 0x99)

TOKEN = re.compile(r"(`[^`]+`)|(\*\*.+?\*\*)|(?<![\w*])(\*[^*\n]+?\*)(?![\w*])"
                   r"|(\[[^\]]+\]\([^)]+\))")


# ---------- markdown -> (title, [sections]) ----------

def cells(row: str) -> list[str]:
    row = row.strip()
    if row.startswith("|"):
        row = row[1:]
    if row.endswith("|") and not row.endswith("\\|"):
        row = row[:-1]
    return [c.strip().replace("\\|", "|") for c in re.split(r"(?<!\\)\|", row)]


def parse(md: str):
    """Returns (doc_title, subtitle, sections); a section is (title, [blocks]).
    Blocks: ('para', text) ('list', [items]) ('code', [lines]) ('table', [rows])."""
    lines = md.split("\n")
    doc_title, subtitle = "", ""
    sections: list[tuple[str, list]] = []
    blocks: list = sections and sections[-1][1]
    para: list[str] = []

    def sink() -> list:
        if not sections:
            sections.append((doc_title or "Document", []))
        return sections[-1][1]

    def flush():
        if para:
            sink().append(("para", " ".join(para)))
            para.clear()

    i = 0
    while i < len(lines):
        line = lines[i]
        if m := re.match(r"(#{1,4})\s+(.*)", line):
            flush()
            if len(m.group(1)) == 1 and not doc_title:
                doc_title = m.group(2)
                while i + 1 < len(lines) and not lines[i + 1].strip():
                    i += 1
                sub = []
                while i + 1 < len(lines) and lines[i + 1].strip() \
                        and not lines[i + 1].startswith(("#", "|", "-", "```")):
                    i += 1
                    sub.append(lines[i].strip())
                subtitle = " ".join(sub)
            else:
                sections.append((m.group(2), []))
        elif line.startswith("```"):
            flush()
            block = []
            i += 1
            while i < len(lines) and not lines[i].startswith("```"):
                block.append(lines[i])
                i += 1
            sink().append(("code", block))
        elif line.lstrip().startswith("|"):
            flush()
            rows = []
            while i < len(lines) and lines[i].lstrip().startswith("|"):
                rows.append(cells(lines[i]))
                i += 1
            i -= 1
            rows = [r for r in rows if not all(re.fullmatch(r":?-+:?", c) for c in r)]
            ncols = len(rows[0])  # normalize ragged rows to the header width
            rows = [r[:ncols - 1] + [" | ".join(r[ncols - 1:])] if len(r) > ncols
                    else r + [""] * (ncols - len(r)) for r in rows]
            sink().append(("table", rows))
        elif re.match(r"(?:-|\d+\.)\s+", line):
            flush()
            items = []
            while i < len(lines) and (m := re.match(r"(-|\d+\.)\s+(.*)", lines[i])):
                text = [("" if m.group(1) == "-" else m.group(1) + " ") + m.group(2)]
                while i + 1 < len(lines) and re.match(r"\s+\S", lines[i + 1]) \
                        and not re.match(r"\s*(?:-|\d+\.)\s", lines[i + 1]):
                    i += 1
                    text.append(lines[i].strip())
                items.append(" ".join(text))
                i += 1
            i -= 1
            sink().append(("list", items))
        elif not line.strip():
            flush()
        else:
            para.append(line.strip())
        i += 1
    flush()
    return doc_title, subtitle, sections


def plain(text: str) -> str:
    return TOKEN.sub(lambda m: re.sub(r"^\W+|\W+$", "", m.group(0).split("](")[0]), text)


# ---------- layout estimation (inches) ----------

def est(block) -> float:
    kind, data = block
    if kind == "para":
        return 0.26 * max(1, -(-len(data) // 92)) + 0.10
    if kind == "list":
        return sum(0.26 * max(1, -(-len(it) // 88)) + 0.06 for it in data) + 0.08
    if kind == "code":
        return 0.17 * len(data) + 0.30
    widths = col_weights(data)
    return 0.30 + sum(row_est(r, widths) for r in data[1:])


def col_weights(rows) -> list[float]:
    w = [min(48, max(6, max(len(r[c]) for r in rows if c < len(r))))
         for c in range(len(rows[0]))]
    return [x / sum(w) for x in w]


def row_est(row, weights) -> float:
    lines = max(max(1, -(-len(c) // max(10, int(f * 130))))
                for c, f in zip(row, weights))
    return 0.235 * lines + 0.10


# ---------- pptx emission ----------

def add_runs(p, text: str, size: int):
    pos = 0
    for m in TOKEN.finditer(text):
        for seg, style in [(text[pos:m.start()], None), (m.group(0), "tok")]:
            if not seg:
                continue
            r = p.add_run()
            r.font.size = Pt(size)
            r.font.color.rgb = INK
            if style is None:
                r.text = seg
            elif seg.startswith("`"):
                r.text = seg[1:-1]
                r.font.name = "Consolas"
                r.font.color.rgb = ACCENT
            elif seg.startswith("**"):
                r.text = seg[2:-2]
                r.font.bold = True
            elif seg.startswith("*"):
                r.text = seg[1:-1]
                r.font.italic = True
            else:
                mm = re.match(r"\[([^\]]+)\]\(([^)]+)\)", seg)
                r.text = mm.group(1)
                r.hyperlink.address = mm.group(2)
        pos = m.end()
    if pos < len(text):
        r = p.add_run()
        r.text = text[pos:]
        r.font.size = Pt(size)
        r.font.color.rgb = INK


class Deck:
    def __init__(self, source: str):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.source = source
        self.slide = None
        self.y = CONTENT_TOP
        self.notes: list[str] = []

    def new_slide(self, title: str):
        self._flush_notes()
        self.slide = self.prs.slides.add_slide(self.blank)
        tb = self.slide.shapes.add_textbox(
            Inches(MARGIN), Inches(0.30), Inches(SLIDE_W - 2 * MARGIN), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        add_runs(p, title, 24)
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = ACCENT
        foot = self.slide.shapes.add_textbox(
            Inches(MARGIN), Inches(SLIDE_H - 0.42), Inches(SLIDE_W - 2 * MARGIN), Inches(0.3))
        fp = foot.text_frame.paragraphs[0]
        fp.text = f"{self.source}   ·   slide {len(self.prs.slides)}"
        fp.runs[0].font.size = Pt(9)
        fp.runs[0].font.color.rgb = DIM
        self.y = CONTENT_TOP

    def _flush_notes(self):
        if self.slide is not None and self.notes:
            self.slide.notes_slide.notes_text_frame.text = "\n".join(self.notes)
        self.notes = []

    def box(self, height: float) -> tuple:
        pos = (Inches(MARGIN), Inches(self.y), Inches(SLIDE_W - 2 * MARGIN), Inches(height))
        self.y += height
        return pos

    def para(self, text: str, size=14, gap=0.10):
        tb = self.slide.shapes.add_textbox(*self.box(est(("para", text)) - gap + 0.02))
        tb.text_frame.word_wrap = True
        add_runs(tb.text_frame.paragraphs[0], text, size)
        self.y += gap
        self.notes.append(plain(text))

    def bullets(self, items: list[str]):
        tb = self.slide.shapes.add_textbox(*self.box(est(("list", items))))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, item in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            add_runs(p, ("" if re.match(r"\d+\. ", item) else "•  ") + item, 13)
            p.space_after = Pt(4)
            self.notes.append("- " + plain(item))

    def code(self, code_lines: list[str]):
        tb = self.slide.shapes.add_textbox(*self.box(est(("code", code_lines))))
        tb.fill.solid()
        tb.fill.fore_color.rgb = CODE_BG
        tb.line.color.rgb = BORDER
        tb.line.width = Pt(0.75)
        tf = tb.text_frame
        tf.word_wrap = False
        for k, ln in enumerate(code_lines):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = ln if ln else " "
            r.font.name = "Consolas"
            r.font.size = Pt(10.5)
            r.font.color.rgb = INK
            p.space_after = Pt(0)

    def table(self, rows: list[list[str]]):
        weights = col_weights(rows)
        height = est(("table", rows))
        shape = self.slide.shapes.add_table(
            len(rows), len(rows[0]), *self.box(height))
        tbl = shape.table
        total = Inches(SLIDE_W - 2 * MARGIN)
        for c, f in enumerate(weights):
            tbl.columns[c].width = int(total * f)
        for ri, row in enumerate(rows):
            tbl.rows[ri].height = Inches(row_est(row, weights) if ri else 0.30)
            for ci, cell_text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG if ri == 0 else RGBColor(255, 255, 255)
                for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                    setattr(cell, m, Pt(4))
                p = cell.text_frame.paragraphs[0]
                add_runs(p, cell_text, 10)
                if ri == 0:
                    for r in p.runs:
                        r.font.bold = True
            self.notes.append(" | ".join(plain(c) for c in row))


def split_block(block, room: float):
    """Largest leading piece of `block` fitting in `room`, and the remainder."""
    kind, data = block
    if kind == "code":
        n = max(3, int((room - 0.30) / 0.17))
        return ("code", data[:n]), (("code", data[n:]) if len(data) > n else None)
    if kind == "list":
        head, used = [], 0.08
        for it in data:
            used += 0.26 * max(1, -(-len(it) // 88)) + 0.06
            if used > room and head:
                break
            head.append(it)
        rest = data[len(head):]
        return ("list", head), (("list", rest) if rest else None)
    if kind == "table":
        weights = col_weights(data)
        head, used = [data[0]], 0.30
        for r in data[1:]:
            used += row_est(r, weights)
            if used > room and len(head) > 1:
                break
            head.append(r)
        rest = data[len(head):]
        return ("table", head), (("table", [data[0]] + rest) if rest else None)
    return block, None


def build(md_path: Path, out_path: Path) -> int:
    doc_title, subtitle, sections = parse(md_path.read_text(encoding="utf-8"))
    deck = Deck(md_path.name)

    title_slide = deck.prs.slides.add_slide(deck.blank)
    tb = title_slide.shapes.add_textbox(
        Inches(MARGIN), Inches(2.5), Inches(SLIDE_W - 2 * MARGIN), Inches(1.2))
    add_runs(tb.text_frame.paragraphs[0], doc_title, 34)
    for r in tb.text_frame.paragraphs[0].runs:
        r.font.bold = True
        r.font.color.rgb = ACCENT
    if subtitle:
        sub = title_slide.shapes.add_textbox(
            Inches(MARGIN), Inches(3.8), Inches(SLIDE_W - 2 * MARGIN), Inches(1.6))
        sub.text_frame.word_wrap = True
        add_runs(sub.text_frame.paragraphs[0], subtitle, 15)
        for r in sub.text_frame.paragraphs[0].runs:
            r.font.color.rgb = DIM

    for title, blocks in sections:
        if not blocks:
            continue
        deck.new_slide(title)
        queue = list(blocks)
        while queue:
            block = queue.pop(0)
            room = CONTENT_BOTTOM - deck.y
            if est(block) > room:
                # Prefer a fresh slide over splitting whenever the block fits
                # whole on one (keeps the §3.4 struct tree in a single box).
                if deck.y > CONTENT_TOP and (est(block) <= BUDGET or room < BUDGET * 0.35):
                    deck.new_slide(title + "  (cont.)")
                    room = BUDGET
                if est(block) > room:
                    head, rest = split_block(block, room)
                    if rest is not None:
                        queue.insert(0, rest)
                    block = head
            kind, data = block
            if kind == "para":
                deck.para(data)
            elif kind == "list":
                deck.bullets(data)
            elif kind == "code":
                deck.code(data)
            else:
                deck.table(data)
    deck._flush_notes()

    deck.prs.core_properties.title = doc_title or md_path.name
    deck.prs.core_properties.comments = "generated by scripts/md2pptx.py — do not commit"
    deck.prs.save(str(out_path))
    return len(deck.prs.slides)


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    md = Path(sys.argv[1]) if len(sys.argv) > 1 else root / "docs" / "ICD-NN.md"
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else md.with_suffix(".pptx")
    n = build(md, out)
    print(f"wrote {out} ({n} slides) — regenerable artifact, do not commit")


if __name__ == "__main__":
    main()
